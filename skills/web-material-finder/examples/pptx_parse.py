# -*- coding: utf-8 -*-
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
path = r'C:\Users\arina\.openclaw\workspace\materials\_inbox\test_geometry_cube.pptx'
prs = Presentation(path)

chunks = []
for i, slide in enumerate(prs.slides, 1):
    parts = []
    if slide.shapes.title:
        parts.append(f"### Slide {i}: {slide.shapes.title.text}")
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    parts.append(t)
    chunks.append("\n".join(parts))

text = "\n\n---\n\n".join(chunks)
print("=== PARSED TEXT (preview) ===")
print(text[:1500])
print("=== TOTAL CHARS:", len(text), "===")
print("=== SLIDES:", len(prs.slides), "===")
