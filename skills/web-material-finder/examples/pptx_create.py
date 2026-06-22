from pptx import Presentation

prs = Presentation()

# Slide 1: title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = 'Задачи на стереометрию: куб'
slide.placeholders[1].text = 'Сборник задач для ЕГЭ профиль, задача 14'

# Slide 2: theory
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Теория: куб ABCDA1B1C1D1'
tf = slide.placeholders[1].text_frame
tf.text = 'Куб — правильный многогранник, все 12 рёбер равны, все углы прямые.'
p = tf.add_paragraph(); p.text = 'Объём куба: V = a^3'
p = tf.add_paragraph(); p.text = 'Площадь поверхности: S = 6a^2'
p = tf.add_paragraph(); p.text = 'Диагональ грани: a*sqrt(2)'
p = tf.add_paragraph(); p.text = 'Диагональ куба: a*sqrt(3)'

# Slide 3: problem 1
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Задача 1 (угол в кубе)'
tf = slide.placeholders[1].text_frame
tf.text = 'В кубе ABCDA1B1C1D1 найдите угол между прямыми AB1 и BC1.'
p = tf.add_paragraph(); p.text = 'Решение: используем векторный метод.'
p = tf.add_paragraph(); p.text = 'Ответ: arccos(1/3) ≈ 70.5 градусов'

# Slide 4: problem 2
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Задача 2 (сечение куба)'
tf = slide.placeholders[1].text_frame
tf.text = 'В кубе с ребром 6 проведено сечение через середины рёбер AB, BC, CC1.'
p = tf.add_paragraph(); p.text = 'Найдите площадь сечения.'
p = tf.add_paragraph(); p.text = 'Ответ: 9*sqrt(3)'

# Slide 5: problem 3 (about cylinder, not cube)
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Задача 3 (цилиндр)'
tf = slide.placeholders[1].text_frame
tf.text = 'В цилиндр с радиусом основания 4 вписана сфера.'
p = tf.add_paragraph(); p.text = 'Найдите объём цилиндра.'
p = tf.add_paragraph(); p.text = 'Ответ: 128π/3'

# Slide 6: parameters with module
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Задача 4 (параметры с модулем)'
tf = slide.placeholders[1].text_frame
tf.text = 'При каких значениях параметра a уравнение |x^2 - 4x + 3| = a имеет ровно 3 решения?'
p = tf.add_paragraph(); p.text = 'Решение: графический метод, модуль даёт симметрию относительно оси x.'
p = tf.add_paragraph(); p.text = 'Ответ: a принадлежит (0; 1) объединение {2}'

# Slide 7: extra
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Дополнительно'
tf = slide.placeholders[1].text_frame
tf.text = 'Полезные формулы для куба: V = a^3, S_полн = 6a^2.'
p = tf.add_paragraph(); p.text = 'Углы в кубе: cos(угол между диагоналями) = 1/3'

prs.save(r'C:\Users\arina\.openclaw\workspace\materials\_inbox\test_geometry_cube.pptx')
print('Saved OK')

prs2 = Presentation(r'C:\Users\arina\.openclaw\workspace\materials\_inbox\test_geometry_cube.pptx')
print('Slides:', len(prs2.slides))
for i, s in enumerate(prs2.slides, 1):
    title = s.shapes.title.text if s.shapes.title else '?'
    print(f'  {i}. {title}')
